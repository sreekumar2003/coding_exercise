"""
Text extraction utilities for PPT, PDF, and plain text inputs.
"""

import io
import re
from pathlib import Path


def extract_from_pptx(file_bytes: bytes) -> str:
    """Extract all text from a .pptx file, preserving slide structure."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
        if parts:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(parts))
    return "\n\n".join(slides_text)


def extract_from_pdf(file_bytes: bytes) -> str:
    """Extract all text from a PDF file."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        text = text.strip()
        if text:
            pages.append(f"[Page {page_num}]\n{text}")
    return "\n\n".join(pages)


def extract_from_text(text: str) -> str:
    """Clean and normalise plain text input."""
    # Collapse excessive blank lines to at most two
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text(file_bytes: bytes | None, filename: str | None, plain_text: str | None) -> str:
    """
    Unified entry point.  Returns raw extracted text regardless of source.
    Raises ValueError when the input cannot be processed.
    """
    if file_bytes and filename:
        ext = Path(filename).suffix.lower()
        if ext in (".pptx", ".ppt"):
            return extract_from_pptx(file_bytes)
        if ext == ".pdf":
            return extract_from_pdf(file_bytes)
        if ext in (".txt", ".md"):
            return extract_from_text(file_bytes.decode("utf-8", errors="replace"))
        raise ValueError(f"Unsupported file type: {ext}")

    if plain_text and plain_text.strip():
        return extract_from_text(plain_text)

    raise ValueError("No valid input provided. Please upload a file or enter text.")
