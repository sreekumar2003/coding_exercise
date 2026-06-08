"""
Modern PowerPoint creator.
Converts a list of Slide dicts into a polished .pptx file and returns the bytes.
"""

import io
from typing import Literal

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Design tokens ─────────────────────────────────────────────────────────────

THEMES: dict[str, dict] = {
    "dark": {
        "bg":         RGBColor(0x0D, 0x1B, 0x2A),   # deep navy
        "accent":     RGBColor(0x00, 0xB4, 0xD8),   # vivid cyan
        "accent2":    RGBColor(0x90, 0xE0, 0xEF),   # light cyan
        "title_fg":   RGBColor(0xFF, 0xFF, 0xFF),
        "heading_fg": RGBColor(0x00, 0xB4, 0xD8),
        "body_fg":    RGBColor(0xE0, 0xE8, 0xF0),
        "bullet_dot": RGBColor(0x00, 0xB4, 0xD8),
        "strip_bg":   RGBColor(0x00, 0xB4, 0xD8),
        "strip_fg":   RGBColor(0xFF, 0xFF, 0xFF),
    },
    "light": {
        "bg":         RGBColor(0xF5, 0xF7, 0xFF),
        "accent":     RGBColor(0x36, 0x5E, 0xF5),
        "accent2":    RGBColor(0x6B, 0x8C, 0xFF),
        "title_fg":   RGBColor(0x1A, 0x1A, 0x2E),
        "heading_fg": RGBColor(0x36, 0x5E, 0xF5),
        "body_fg":    RGBColor(0x2D, 0x2D, 0x3D),
        "bullet_dot": RGBColor(0x36, 0x5E, 0xF5),
        "strip_bg":   RGBColor(0x36, 0x5E, 0xF5),
        "strip_fg":   RGBColor(0xFF, 0xFF, 0xFF),
    },
    "minimal": {
        "bg":         RGBColor(0xFF, 0xFF, 0xFF),
        "accent":     RGBColor(0xFF, 0x6B, 0x35),
        "accent2":    RGBColor(0xFF, 0xA0, 0x70),
        "title_fg":   RGBColor(0x1C, 0x1C, 0x1C),
        "heading_fg": RGBColor(0xFF, 0x6B, 0x35),
        "body_fg":    RGBColor(0x33, 0x33, 0x33),
        "bullet_dot": RGBColor(0xFF, 0x6B, 0x35),
        "strip_bg":   RGBColor(0xFF, 0x6B, 0x35),
        "strip_fg":   RGBColor(0xFF, 0xFF, 0xFF),
    },
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── low-level drawing helpers ─────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor, alpha: int = 255):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()   # no border
    return shape


def _add_textbox(slide, left, top, width, height,
                 text: str, font_size: int, color: RGBColor,
                 bold=False, align=PP_ALIGN.LEFT, italic=False) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _add_bullet_para(tf, text: str, font_size: int, body_color: RGBColor, bullet_color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree

    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    para.level = 0

    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = body_color

    # Custom bullet: colored ▶ via XML
    pPr = para._p.get_or_add_pPr()
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    srgbClr = etree.SubElement(buClr, qn("a:srgbClr"))
    srgbClr.set("val", str(bullet_color))
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "▶")
    buSzPct = etree.SubElement(pPr, qn("a:buSzPct"))
    buSzPct.set("val", "70000")   # 70 %


# ── slide builders ────────────────────────────────────────────────────────────

def _build_title_slide(slide, slide_data: dict, theme: dict):
    _set_bg(slide, theme["bg"])

    # Left accent bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, theme["accent"])

    # Bottom strip
    _add_rect(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0), theme["strip_bg"])

    # Main title
    _add_textbox(
        slide,
        Inches(0.5), Inches(2.0), Inches(12.0), Inches(1.8),
        slide_data["title"], font_size=44,
        color=theme["title_fg"], bold=True, align=PP_ALIGN.LEFT,
    )

    # Subtitle bullets
    if slide_data["bullets"]:
        subtitle = "  •  ".join(slide_data["bullets"])
        _add_textbox(
            slide,
            Inches(0.5), Inches(3.9), Inches(11.0), Inches(0.8),
            subtitle, font_size=18,
            color=theme["accent2"], bold=False, align=PP_ALIGN.LEFT,
        )

    # Bottom strip label
    _add_textbox(
        slide,
        Inches(0.3), Inches(6.55), Inches(8.0), Inches(0.6),
        "Generated with Modern PPT Creator", font_size=12,
        color=theme["strip_fg"], align=PP_ALIGN.LEFT,
    )


def _build_section_slide(slide, slide_data: dict, theme: dict):
    _set_bg(slide, theme["bg"])

    # Full-width accent bar top third
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.8), theme["accent"])

    # Title in accent area
    _add_textbox(
        slide,
        Inches(0.6), Inches(0.8), Inches(12.0), Inches(1.4),
        slide_data["title"], font_size=40,
        color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.LEFT,
    )

    # Sub-text / single bullet
    if slide_data["bullets"]:
        _add_textbox(
            slide,
            Inches(0.6), Inches(3.2), Inches(12.0), Inches(0.7),
            slide_data["bullets"][0], font_size=22,
            color=theme["body_fg"], italic=True,
        )


def _build_content_slide(slide, slide_data: dict, theme: dict, slide_number: int):
    _set_bg(slide, theme["bg"])

    # Thin top accent bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), theme["accent"])

    # Side accent strip
    _add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, theme["accent"])

    # Slide number badge (top-right)
    _add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.45), theme["accent"])
    _add_textbox(
        slide,
        Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.45),
        str(slide_number), font_size=14,
        color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.CENTER,
    )

    # Title
    _add_textbox(
        slide,
        Inches(0.3), Inches(0.3), Inches(11.8), Inches(0.9),
        slide_data["title"], font_size=32,
        color=theme["heading_fg"], bold=True,
    )

    # Underline rule
    _add_rect(slide, Inches(0.3), Inches(1.25), Inches(11.8), Inches(0.04), theme["accent"])

    # Bullet points
    if slide_data["bullets"]:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        # Remove default first paragraph
        tf.paragraphs[0].text = ""

        for bullet in slide_data["bullets"]:
            _add_bullet_para(tf, bullet, font_size=20, body_color=theme["body_fg"],
                             bullet_color=theme["bullet_dot"])


# ── public API ────────────────────────────────────────────────────────────────

def create_pptx(slides_data: list[dict], theme_name: str = "dark") -> bytes:
    """
    Build a modern PPTX from *slides_data* and return the raw bytes.

    Parameters
    ----------
    slides_data : list of Slide TypedDicts  {title, bullets, layout}
    theme_name  : one of "dark" | "light" | "minimal"
    """
    theme = THEMES.get(theme_name, THEMES["dark"])

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout for full design control
    blank_layout = prs.slide_layouts[6]

    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_data.get("layout", "content")

        if layout == "title":
            _build_title_slide(slide, slide_data, theme)
        elif layout == "section":
            _build_section_slide(slide, slide_data, theme)
        else:
            _build_content_slide(slide, slide_data, theme, slide_number=idx)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
